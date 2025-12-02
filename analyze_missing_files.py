#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Script d'analyse des fichiers Excel manquants pour le cahier des charges
"""

import sys
import io

# Forcer UTF-8 pour l'encodage
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

try:
    import openpyxl
except ImportError:
    print("openpyxl n'est pas installé. Installation en cours...")
    import subprocess
    subprocess.check_call([sys.executable, "-m", "pip", "install", "openpyxl"])
    import openpyxl

from pathlib import Path
import json

def analyze_excel_kpi(file_path):
    """Analyse un fichier Excel de type KPI/Dashboard"""
    print(f"\n{'='*80}")
    print(f"Analyse du fichier: {file_path.name}")
    print(f"{'='*80}\n")
    
    try:
        wb = openpyxl.load_workbook(file_path, data_only=True)
        sheet_names = wb.sheetnames
        print(f"Nombre de feuilles: {len(sheet_names)}")
        print(f"Noms des feuilles: {sheet_names}\n")
        
        analysis = {
            'file_name': file_path.name,
            'file_path': str(file_path),
            'number_of_sheets': len(sheet_names),
            'sheets': [],
            'type': 'KPI/Dashboard'
        }
        
        for sheet_name in sheet_names:
            print(f"\n--- Feuille: {sheet_name} ---")
            
            try:
                ws = wb[sheet_name]
                max_row = ws.max_row
                max_col = ws.max_column
                
                print(f"Dimensions réelles: {max_row} lignes x {max_col} colonnes")
                
                # Analyser les premières lignes pour détecter la structure
                header_candidates = []
                sample_data = []
                kpi_indicators = []
                
                for row_idx in range(1, min(30, max_row + 1)):
                    row_values = []
                    non_null_count = 0
                    
                    for col_idx in range(1, min(max_col + 1, 50)):
                        cell = ws.cell(row=row_idx, column=col_idx)
                        value = cell.value
                        
                        if value is not None:
                            try:
                                str_value = str(value).strip()
                                if str_value:
                                    non_null_count += 1
                                    row_values.append(str_value[:80])
                                    # Détecter les indicateurs KPI
                                    if any(keyword in str_value.upper() for keyword in ['KPI', 'RATIO', 'PERFORMANCE', 'METRIC', 'TARGET', 'ACTUAL']):
                                        kpi_indicators.append(str_value)
                            except:
                                pass
                        else:
                            row_values.append("")
                    
                    if non_null_count > max_col * 0.2:
                        header_candidates.append(row_idx - 1)
                        if len(sample_data) < 15:
                            sample_data.append({
                                'row': row_idx,
                                'values': row_values[:20]
                            })
                
                if header_candidates:
                    print(f"\nCandidats pour les en-têtes (lignes): {header_candidates[:10]}")
                    
                    print(f"\nPremières lignes de données (échantillon):")
                    for sample in sample_data[:8]:
                        print(f"  Ligne {sample['row']}: {sample['values'][:8]}")
                
                if kpi_indicators:
                    print(f"\nIndicateurs KPI détectés: {kpi_indicators[:10]}")
                
                # Détecter les cellules fusionnées
                merged_cells = []
                try:
                    if hasattr(ws, 'merged_cells') and ws.merged_cells:
                        merged_cells = list(ws.merged_cells.ranges)
                except AttributeError:
                    pass
                
                if merged_cells:
                    print(f"\nCellules fusionnées détectées: {len(merged_cells)}")
                    print(f"Premières cellules fusionnées: {[str(mc) for mc in merged_cells[:5]]}")
                
                # Détecter les formules
                formula_count = 0
                formula_samples = []
                for row in ws.iter_rows(max_row=min(200, max_row), max_col=min(max_col, 50)):
                    for cell in row:
                        if cell.data_type == 'f':
                            formula_count += 1
                            if len(formula_samples) < 5:
                                try:
                                    formula_samples.append({
                                        'cell': cell.coordinate,
                                        'formula': str(cell.value)[:100]
                                    })
                                except:
                                    pass
                
                if formula_count > 0:
                    print(f"\nFormules détectées (dans les 200 premières lignes): {formula_count}")
                    if formula_samples:
                        print("Exemples de formules:")
                        for f in formula_samples[:3]:
                            print(f"  {f['cell']}: {f['formula'][:60]}...")
                
                sheet_analysis = {
                    'name': sheet_name,
                    'max_rows': max_row,
                    'max_columns': max_col,
                    'merged_cells_count': len(merged_cells),
                    'has_formulas': formula_count > 0,
                    'formula_count_sample': formula_count,
                    'header_candidates': header_candidates[:10],
                    'sample_data': sample_data[:10],
                    'kpi_indicators': kpi_indicators[:20]
                }
                analysis['sheets'].append(sheet_analysis)
                
            except Exception as e:
                print(f"Erreur lors de l'analyse de la feuille {sheet_name}: {str(e)}")
                analysis['sheets'].append({
                    'name': sheet_name,
                    'error': str(e)
                })
        
        wb.close()
        return analysis
        
    except Exception as e:
        print(f"Erreur lors de l'analyse du fichier {file_path}: {str(e)}")
        return {
            'file_name': file_path.name,
            'error': str(e)
        }

def analyze_pptx(file_path):
    """Analyse un fichier PowerPoint - tentative basique"""
    print(f"\n{'='*80}")
    print(f"Analyse du fichier PowerPoint: {file_path.name}")
    print(f"{'='*80}\n")
    
    try:
        try:
            from pptx import Presentation
        except ImportError:
            print("python-pptx n'est pas installé. Tentative d'installation...")
            import subprocess
            subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
            from pptx import Presentation
        
        prs = Presentation(file_path)
        
        analysis = {
            'file_name': file_path.name,
            'file_path': str(file_path),
            'number_of_slides': len(prs.slides),
            'type': 'PowerPoint Presentation',
            'slides': []
        }
        
        print(f"Nombre de diapositives: {len(prs.slides)}\n")
        
        for idx, slide in enumerate(prs.slides, 1):
            print(f"--- Diapositive {idx} ---")
            
            slide_content = {
                'slide_number': idx,
                'title': '',
                'text_content': [],
                'tables': 0,
                'charts': 0
            }
            
            # Extraire le texte
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text.strip()
                    if text:
                        if shape.shape_type == 1:  # AutoShape with text
                            if not slide_content['title']:
                                slide_content['title'] = text.split('\n')[0][:100]
                            slide_content['text_content'].append(text[:200])
                        elif shape.shape_type == 5:  # Group
                            slide_content['text_content'].append(text[:200])
                
                # Compter les tableaux
                if shape.shape_type == 19:  # Table
                    slide_content['tables'] += 1
                
                # Compter les graphiques
                if shape.shape_type == 3:  # Chart
                    slide_content['charts'] += 1
            
            if slide_content['text_content']:
                print(f"  Titre: {slide_content['title']}")
                print(f"  Nombre de textes: {len(slide_content['text_content'])}")
                print(f"  Tableaux: {slide_content['tables']}")
                print(f"  Graphiques: {slide_content['charts']}")
                print(f"  Échantillon de contenu:")
                for text in slide_content['text_content'][:3]:
                    print(f"    - {text[:100]}")
            
            analysis['slides'].append(slide_content)
        
        return analysis
        
    except ImportError:
        print("Impossible d'installer python-pptx. Analyse basique du fichier...")
        return {
            'file_name': file_path.name,
            'type': 'PowerPoint Presentation',
            'note': 'L\'analyse détaillée nécessite python-pptx'
        }
    except Exception as e:
        print(f"Erreur lors de l'analyse du fichier PowerPoint {file_path}: {str(e)}")
        return {
            'file_name': file_path.name,
            'error': str(e)
        }

def main():
    base_path = Path("docs")
    
    # Fichiers à analyser
    files_to_analyze = [
        (base_path / "Utility Folder" / "Utility-Ratio.xlsx", "excel"),
        (base_path / "Utility Folder" / "UTILITY KPI DASH BOARD FOR EXCELLENCE.xlsx", "excel"),
        (base_path / "RE_ Follow-up on the digitalization of your department's needs – Coordination and next steps" / "Maintenance KPI.pptx", "pptx"),
    ]
    
    all_analyses = []
    
    for file_path, file_type in files_to_analyze:
        if file_path.exists():
            if file_type == "excel":
                analysis = analyze_excel_kpi(file_path)
                all_analyses.append(analysis)
            elif file_type == "pptx":
                analysis = analyze_pptx(file_path)
                all_analyses.append(analysis)
        else:
            print(f"Fichier non trouvé: {file_path}")
    
    # Sauvegarder l'analyse dans un fichier JSON
    output_file = Path("missing_files_analysis.json")
    with open(output_file, 'w', encoding='utf-8') as f:
        json.dump(all_analyses, f, ensure_ascii=False, indent=2)
    
    print(f"\n\n{'='*80}")
    print("Analyse complète terminée. Résultats sauvegardés dans missing_files_analysis.json")
    print(f"{'='*80}")

if __name__ == "__main__":
    main()

